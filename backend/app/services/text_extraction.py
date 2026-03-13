"""Extract text content from PDF, PPTX, and DOCX files."""

import io
import logging

logger = logging.getLogger(__name__)

MAX_EXTRACTED_CHARS = 100_000

# Magic byte signatures for file validation
MAGIC_BYTES = {
    "pdf": b"%PDF",
    "pptx": b"PK",  # ZIP-based (Office Open XML)
    "docx": b"PK",  # ZIP-based (Office Open XML)
}


def _validate_magic_bytes(file_bytes: bytes, file_type: str) -> bool:
    """Check that file content matches expected magic bytes."""
    expected = MAGIC_BYTES.get(file_type)
    if not expected:
        return False
    return file_bytes[:len(expected)] == expected


def _extract_pdf(file_bytes: bytes) -> str | None:
    try:
        import pdfplumber

        text_parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n\n".join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("PDF text extraction failed: %s", e)
        return None


def _extract_pptx(file_bytes: bytes) -> str | None:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("PPTX text extraction failed: %s", e)
        return None


def _extract_docx(file_bytes: bytes) -> str | None:
    try:
        import docx

        doc = docx.Document(io.BytesIO(file_bytes))
        text_parts = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                text_parts.append(text)
        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    text_parts.append(row_text)
        return "\n\n".join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("DOCX text extraction failed: %s", e)
        return None


def extract_text(file_bytes: bytes, file_type: str) -> str | None:
    """Extract text from a file. Returns None if extraction fails.

    Args:
        file_bytes: Raw file content
        file_type: File extension without dot (e.g. "pdf", "pptx", "docx")
    """
    file_type = file_type.lower().lstrip(".")

    if not _validate_magic_bytes(file_bytes, file_type):
        logger.warning("File magic bytes don't match expected type: %s", file_type)
        return None

    extractors = {
        "pdf": _extract_pdf,
        "pptx": _extract_pptx,
        "docx": _extract_docx,
    }

    extractor = extractors.get(file_type)
    if not extractor:
        logger.warning("No extractor for file type: %s", file_type)
        return None

    result = extractor(file_bytes)
    if result and len(result) > MAX_EXTRACTED_CHARS:
        result = result[:MAX_EXTRACTED_CHARS] + "\n\n[... text truncated at 100K characters]"
    return result
